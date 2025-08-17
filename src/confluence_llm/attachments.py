from __future__ import annotations

import io
import logging
import re
from typing import Any

"""
Attachment indexing helpers.
- Extracts text from PDF, DOCX, PPTX, PPT (legacy), XLSX, XLS (legacy), TXT/CSV/MD.
- All imports are optional; if a format library isn't present, we skip gracefully.
"""

# Optional third-party imports (guarded)
try:
    import PyPDF2 as _PyPDF2
except Exception:  # pragma: no cover
    _PyPDF2 = None


try:
    import docx as _docx  # python-docx
except Exception:  # pragma: no cover
    _docx = None

try:
    from pptx import Presentation as _Presentation  # python-pptx
except Exception:  # pragma: no cover
    _Presentation = None

try:
    import openpyxl as _openpyxl  # for .xlsx
except Exception:  # pragma: no cover
    _openpyxl = None

try:
    import xlrd as _xlrd  # for .xls (use xlrd<2.0 for legacy support)
except Exception:  # pragma: no cover
    _xlrd = None

try:
    import olefile as _olefile  # for legacy .ppt OLE parsing
except Exception:  # pragma: no cover
    _olefile = None

# Expose typed module variables (Any | None so mypy is happy)
PyPDF2: Any | None = _PyPDF2
docx: Any | None = _docx
Presentation: Any | None = _Presentation
openpyxl: Any | None = _openpyxl
xlrd: Any | None = _xlrd
olefile: Any | None = _olefile

logger = logging.getLogger(__name__)


def _extract_strings_from_bytes(data: bytes) -> str:
    """Heuristic: pull ASCII and UTF-16LE strings out of arbitrary bytes (for legacy formats)."""
    # ASCII-ish
    ascii_strings = re.findall(rb"[ -~]{4,}", data)
    # UTF-16LE: pairs of byte + null
    try:
        utf16le = data.decode("utf-16le", errors="ignore")
        utf16_strings = re.findall(r"[\w\s\-\.,:;/%\(\)\[\]\{\}\+\=]{4,}", utf16le)
    except Exception:
        utf16_strings = []
    parts_list = [s.decode("ascii", errors="ignore") for s in ascii_strings] + utf16_strings
    # de-dup while preserving order
    seen: set[str] = set()
    out: list[str] = []
    for p in parts_list:
        p = p.strip()
        if p and p not in seen:
            seen.add(p)
            out.append(p)
    return "\n".join(out)


def sniff_text(filename: str, data: bytes) -> str:
    name = (filename or "").lower()

    # PDF
    if name.endswith(".pdf") and PyPDF2:
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(data))
            texts: list[str] = []
            for page in reader.pages:
                try:
                    page_text = page.extract_text() or ""
                except Exception as e:
                    logger.warning("PDF page text extraction failed: %s", e)
                    page_text = ""
                if page_text:
                    texts.append(page_text)
            return "\n".join(texts)
        except Exception as e:
            logger.warning("PDF parsing failed: %s", e)
            return ""

    # Word (docx / doc)
    if (name.endswith(".docx") or name.endswith(".doc")) and docx:
        try:
            d = docx.Document(io.BytesIO(data))
            return "\n".join(p.text for p in d.paragraphs)
        except Exception as e:
            logger.warning("DOCX parsing failed: %s", e)
            return ""

    # PowerPoint (pptx)
    if name.endswith(".pptx") and Presentation:
        try:
            prs = Presentation(io.BytesIO(data))
            ppt_texts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    # Prefer has_text_frame when available; fall back to a best-effort "text" attr.
                    has_text = getattr(shape, "has_text_frame", False)
                    if has_text and getattr(shape, "text_frame", None):
                        ppt_texts.append(str(shape.text_frame.text or ""))
                    else:
                        text_attr = getattr(shape, "text", None)
                        if text_attr:
                            ppt_texts.append(str(text_attr))
            return "\n".join(t for t in ppt_texts if t)
        except Exception as e:
            logger.warning("PPTX parsing failed: %s", e)
            return ""

    # PowerPoint (legacy .ppt via OLE heuristic)
    if name.endswith(".ppt") and olefile:
        try:
            with olefile.OleFileIO(io.BytesIO(data)) as ole:
                ole_texts: list[str] = []
                for stream_name in ole.listdir(streams=True, storages=False):
                    try:
                        with ole.openstream(stream_name) as s:
                            raw = s.read()
                        ole_texts.append(_extract_strings_from_bytes(raw))
                    except Exception as e:
                        logger.debug("OLE stream read failed for %s: %s", stream_name, e)
                return "\n".join([p for p in ole_texts if p])
        except Exception as e:
            logger.warning("PPT (OLE) parse failed, falling back to heuristic: %s", e)
            # Fallback: heuristic across full file if OLE parse fails
            return _extract_strings_from_bytes(data)

    # Excel (xlsx)
    if name.endswith(".xlsx") and openpyxl:
        try:
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            xlsx_rows: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        xlsx_rows.append("\t".join(cells))
            return "\n".join(xlsx_rows)
        except Exception as e:
            logger.warning("XLSX parsing failed: %s", e)
            return ""

    # Excel (legacy .xls)
    if name.endswith(".xls") and xlrd:
        try:
            book = xlrd.open_workbook(file_contents=data)
            xls_rows: list[str] = []
            for si in range(book.nsheets):
                sh = book.sheet_by_index(si)
                for rx in range(sh.nrows):
                    row = sh.row_values(rx)
                    cells = [str(c) for c in row if c not in (None, "")]
                    if cells:
                        xls_rows.append("\t".join(cells))
            return "\n".join(xls_rows)
        except Exception as e:
            logger.warning("XLS parsing failed: %s", e)
            return ""

    # Plaintext-ish
    if name.endswith((".txt", ".csv", ".md")):
        try:
            return data.decode("utf-8", errors="ignore")
        except Exception as e:
            logger.debug("UTF-8 decode failed, ignoring: %s", e)
            return ""

    # Fallback: unknown format
    return ""
